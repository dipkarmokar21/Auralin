# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Premium color palette
C_BG      = RGBColor(0x05, 0x05, 0x10)  # ultra dark navy
C_CARD    = RGBColor(0x10, 0x10, 0x22)  # card dark
C_CARD2   = RGBColor(0x18, 0x18, 0x30)  # slightly lighter card
C_PURPLE  = RGBColor(0x7B, 0x2F, 0xFF)  # vivid purple
C_VIOLET  = RGBColor(0x4A, 0x00, 0xC8)  # deep violet
C_RED     = RGBColor(0xFA, 0x2D, 0x48)  # Auralin red
C_PINK    = RGBColor(0xFF, 0x6B, 0x9D)  # soft pink
C_CYAN    = RGBColor(0x00, 0xD4, 0xFF)  # cyan accent
C_GOLD    = RGBColor(0xFF, 0xC8, 0x00)  # gold
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0x9A, 0x9A, 0xBB)
C_MGRAY   = RGBColor(0x3A, 0x3A, 0x5C)

PICS = [
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (1).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (2).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (3).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (4).png",
]

# ── Helpers ───────────────────────────────────────────────────────────────
def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri Light"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = font
    return tb

def bullets(slide, items, l, t, w, h, base=16):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if item.startswith("##"):
            r = p.add_run(); r.text = item[2:]
            r.font.size = Pt(base+2); r.font.bold = True
            r.font.color.rgb = C_CYAN; r.font.name = "Calibri"
        else:
            lvl = 1 if item.startswith("  ") else 0
            p.level = lvl
            r = p.add_run()
            r.text = ("  " if lvl==0 else "    ") + item.strip()
            r.font.size = Pt(base if lvl==0 else base-1)
            r.font.color.rgb = C_WHITE if lvl==0 else C_LGRAY
            r.font.name = "Calibri Light"
    return tb

def add_pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return None

def transition_none(slide):
    """No auto-advance — manual click only."""
    tr = slide._element.find(qn('p:transition'))
    if tr is not None:
        slide._element.remove(tr)
    tr = etree.SubElement(slide._element, qn('p:transition'))
    tr.set('spd', 'slow')
    # NO advTm attribute = manual advance only
    fade = etree.SubElement(tr, qn('p:fade'))
    fade.set('thruBlk', '0')

def add_anim(slide, shape, delay=0):
    """Simple fade-in animation."""
    try:
        timing = slide._element.find(qn('p:timing'))
        if timing is None:
            timing = etree.SubElement(slide._element, qn('p:timing'))
        tnLst = timing.find(qn('p:tnLst'))
        if tnLst is None:
            tnLst = etree.SubElement(timing, qn('p:tnLst'))
        par = tnLst.find(qn('p:par'))
        if par is None:
            par = etree.SubElement(tnLst, qn('p:par'))
            cTn = etree.SubElement(par, qn('p:cTn'))
            cTn.set('id','1'); cTn.set('dur','indefinite')
            cTn.set('restart','whenNotActive'); cTn.set('nodeType','tmRoot')
            cl = etree.SubElement(cTn, qn('p:childTnLst'))
            seq = etree.SubElement(cl, qn('p:seq'))
            seq.set('concurrent','1'); seq.set('nextAc','seek')
            cTn2 = etree.SubElement(seq, qn('p:cTn'))
            cTn2.set('id','2'); cTn2.set('dur','indefinite')
            cTn2.set('nodeType','mainSeq')
            etree.SubElement(cTn2, qn('p:childTnLst'))
            etree.SubElement(seq, qn('p:prevCondLst'))
            nc = etree.SubElement(seq, qn('p:nextCondLst'))
            cond = etree.SubElement(nc, qn('p:cond'))
            cond.set('evt','onNext'); cond.set('delay','0')
            tEl = etree.SubElement(cond, qn('p:tn'))
            tEl.set('val','2')
        seq = par.find(qn('p:cTn')).find(qn('p:childTnLst')).find(qn('p:seq'))
        mainTnLst = seq.find(qn('p:cTn')).find(qn('p:childTnLst'))
        idx = len(mainTnLst) + 3
        p2 = etree.SubElement(mainTnLst, qn('p:par'))
        cTn3 = etree.SubElement(p2, qn('p:cTn'))
        cTn3.set('id', str(idx)); cTn3.set('fill','hold')
        sc = etree.SubElement(cTn3, qn('p:stCondLst'))
        c = etree.SubElement(sc, qn('p:cond')); c.set('delay', str(delay))
        cl2 = etree.SubElement(cTn3, qn('p:childTnLst'))
        p3 = etree.SubElement(cl2, qn('p:par'))
        cTn4 = etree.SubElement(p3, qn('p:cTn'))
        cTn4.set('id', str(idx+1)); cTn4.set('fill','hold')
        sc2 = etree.SubElement(cTn4, qn('p:stCondLst'))
        c2 = etree.SubElement(sc2, qn('p:cond')); c2.set('delay','0')
        cl3 = etree.SubElement(cTn4, qn('p:childTnLst'))
        ae = etree.SubElement(cl3, qn('p:animEffect'))
        ae.set('transition','in'); ae.set('filter','fade')
        cb = etree.SubElement(ae, qn('p:cBhvr'))
        cTn5 = etree.SubElement(cb, qn('p:cTn'))
        cTn5.set('id', str(idx+2)); cTn5.set('dur','600'); cTn5.set('fill','hold')
        te = etree.SubElement(cb, qn('p:tgtEl'))
        sp = etree.SubElement(te, qn('p:spTgt'))
        sp.set('spid', str(shape.shape_id))
    except Exception:
        pass

def premium_bg(slide):
    """Premium layered background."""
    bg(slide, C_BG)
    # Large dark blobs for depth
    rect(slide, -2, -2, 8, 6, C_VIOLET)
    rect(slide, 7, 3, 8, 6, RGBColor(0x3A,0x00,0x80))
    # Dark overlay to make it subtle
    rect(slide, 0, 0, 13.33, 7.5, C_BG)
    # Subtle glow circles
    rect(slide, 10.5, -1, 4, 4, C_PURPLE)
    rect(slide, 10.7, -0.8, 3.6, 3.6, C_BG)
    rect(slide, -1, 5, 4, 4, C_RED)
    rect(slide, -0.8, 5.2, 3.6, 3.6, C_BG)

def header_bar(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.2, C_CARD)
    rect(slide, 0, 0, 0.22, 1.2, C_PURPLE)
    rect(slide, 0.22, 1.15, 13.11, 0.05, C_VIOLET)
    t = txt(slide, title, 0.45, 0.08, 11.5, 0.75, size=36, bold=True,
            color=C_WHITE, font="Calibri")
    if sub:
        txt(slide, sub, 0.45, 0.82, 11.5, 0.32, size=13,
            color=C_LGRAY, italic=True)
    txt(slide, "♪", 12.3, 0.1, 0.9, 0.9, size=36, bold=True,
        color=C_PURPLE, align=PP_ALIGN.CENTER)
    return t

def footer_bar(slide, n):
    rect(slide, 0, 7.1, 13.33, 0.4, C_CARD)
    rect(slide, 0, 7.1, 0.22, 0.4, C_RED)
    txt(slide, "Auralin Music Player  |  Dip Karmokar  |  OOP Project 2026",
        0.35, 7.13, 10, 0.28, size=11, color=C_LGRAY)
    txt(slide, f"{n} / 15", 12.4, 7.13, 0.85, 0.28,
        size=11, color=C_LGRAY, align=PP_ALIGN.RIGHT)

def card(slide, l, t, w, h, accent=None):
    r = rect(slide, l, t, w, h, C_CARD2)
    if accent:
        rect(slide, l, t, w, 0.06, accent)
    return r

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl)
# Big glow ring
rect(sl, 8.8, 0.3, 4.2, 4.2, C_PURPLE)
rect(sl, 9.1, 0.6, 3.6, 3.6, C_VIOLET)
rect(sl, 9.4, 0.9, 3.0, 3.0, C_BG)
t_note = txt(sl, "♪", 9.7, 1.0, 2.5, 2.5, size=88, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
# Thin accent line left
rect(sl, 0.5, 1.5, 0.08, 3.5, C_RED)
t1 = txt(sl, "AURALIN", 0.8, 1.4, 9, 1.6, size=82, bold=True, color=C_WHITE, font="Calibri")
t2 = txt(sl, "MUSIC PLAYER", 0.8, 2.85, 9, 1.0, size=46, bold=True, color=C_RED, font="Calibri")
rect(sl, 0.8, 3.95, 7, 0.06, C_PURPLE)
t3 = txt(sl, "A Modern JavaFX Desktop Music Application", 0.8, 4.1, 10, 0.55, size=20, color=C_LGRAY, italic=True)
t4 = txt(sl, "Dip Karmokar   |   OOP Project   |   2026", 0.8, 4.75, 10, 0.45, size=16, color=C_LGRAY)
rect(sl, 0, 6.75, 13.33, 0.75, C_CARD)
rect(sl, 0, 6.75, 13.33, 0.06, C_PURPLE)
txt(sl, "JavaFX 25   ·   mp3agic   ·   MVC Architecture   ·   Windows Installer   ·   Bundled JRE",
    0.5, 6.85, 12.33, 0.45, size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)
transition_none(sl)
for s,d in [(t1,0),(t2,200),(t3,400),(t4,600)]: add_anim(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Agenda","What we'll cover today"); footer_bar(sl,2)
transition_none(sl)
items = [
    ("01","Introduction & Motivation",C_RED),
    ("02","System Architecture",C_PURPLE),
    ("03","Key Features",C_CYAN),
    ("04","Implementation & Patterns",C_RED),
    ("05","Data Persistence",C_PURPLE),
    ("06","UI/UX Design",C_CYAN),
    ("07","Challenges & Solutions",C_RED),
    ("08","Conclusion & Future Work",C_PURPLE),
]
for i,(num,label,col) in enumerate(items):
    row=i//2; ci=i%2
    fx=0.4+ci*6.5; fy=1.35+row*1.35
    r=rect(sl,fx,fy,6.1,1.1,C_CARD2)
    rect(sl,fx,fy,0.5,1.1,col)
    txt(sl,num,fx+0.05,fy+0.28,0.45,0.55,size=15,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,label,fx+0.65,fy+0.28,5.3,0.55,size=18,bold=True,color=C_WHITE)
    add_anim(sl,r,i*120)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Introduction","What is Auralin Music Player?"); footer_bar(sl,3)
transition_none(sl)
c1=card(sl,0.4,1.3,6.0,5.5,C_PURPLE); c2=card(sl,6.8,1.3,6.1,5.5,C_RED)
txt(sl,"What is Auralin?",0.6,1.38,5.6,0.45,size=16,bold=True,color=C_CYAN)
txt(sl,"Goals & Scope",7.0,1.38,5.7,0.45,size=16,bold=True,color=C_CYAN)
b1=bullets(sl,["Modern desktop music player","Built with JavaFX 25 + Java","Inspired by Spotify dark UI","##Motivation","Existing players lack modern UI","Deep dive into JavaFX MVC","##Tech Stack","JavaFX 25, mp3agic, Inno Setup","launch4j, jlink bundled JRE"],0.55,1.85,5.7,4.8)
b2=bullets(sl,["Fast startup — no MediaPlayer on load","Persistent library across sessions","Open With — default MP3 player","Custom frameless Win11 window","##Deliverables","Runnable .exe installer","Bundled JRE — no Java needed","Full MVC architecture"],6.95,1.85,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: add_anim(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"System Architecture","MVC + Observer + Facade + Strategy"); footer_bar(sl,4)
transition_none(sl)
arch=[
    (0.3,1.35,3.9,5.4,C_RED,"MODEL",["Song.java","title, artist","filePath","liked, plays","lastPlayed","getDurationStr()"]),
    (4.7,1.35,3.9,5.4,C_PURPLE,"VIEW",["HomeView","LibraryView","SearchView","LikedView","SongTable","PlayerBar","Sidebar","NowPlayingView"]),
    (9.1,1.35,3.9,5.4,C_CYAN,"CONTROLLER",["PlayerController","ViewManager","DatabaseService","FileImportService","PlayerBar.Listener","MusicCard.Listener"]),
]
for bx,by,bw,bh,bc,bt,bi in arch:
    r=card(sl,bx,by,bw,bh,bc)
    txt(sl,bt,bx+0.1,by+0.1,bw-0.2,0.5,size=17,bold=True,color=bc,align=PP_ALIGN.CENTER,font="Calibri")
    for j,item in enumerate(bi):
        txt(sl,"  "+item,bx+0.2,by+0.7+j*0.72,bw-0.3,0.6,size=13,color=C_LGRAY)
    add_anim(sl,r,200)
txt(sl,"<-->",4.15,3.8,0.55,0.5,size=18,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER)
txt(sl,"<-->",8.6,3.8,0.55,0.5,size=18,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER)
txt(sl,"Observer Pattern  |  MVC  |  Facade  |  Strategy  |  Singleton",1.5,6.82,10.33,0.38,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Key Features","Everything Auralin can do"); footer_bar(sl,5)
transition_none(sl)
feats=[
    ("Library","Add files & folders",C_RED),
    ("Search","Real-time filtering",C_PURPLE),
    ("Liked Songs","Heart toggle + persist",C_CYAN),
    ("Home View","Recent + Recommended",C_RED),
    ("Now Playing","Full-screen overlay",C_PURPLE),
    ("Auto-Save","Persists across sessions",C_CYAN),
    ("Open With","Default MP3 player",C_RED),
    ("Custom Window","Frameless + Win11 btns",C_PURPLE),
]
pos=[(0.3,1.3),(3.55,1.3),(6.8,1.3),(10.05,1.3),(0.3,3.9),(3.55,3.9),(6.8,3.9),(10.05,3.9)]
for i,((fx,fy),(title,desc,col)) in enumerate(zip(pos,feats)):
    r=card(sl,fx,fy,3.0,2.3,col)
    rect(sl,fx,fy+1.6,3.0,0.06,col)
    txt(sl,title,fx+0.15,fy+0.2,2.7,0.55,size=16,bold=True,color=col,font="Calibri")
    txt(sl,desc,fx+0.15,fy+0.85,2.7,0.65,size=13,color=C_LGRAY)
    add_anim(sl,r,i*100)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Screenshots 1 & 2
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Home View & Library","Screenshots from Auralin Music Player"); footer_bar(sl,6)
transition_none(sl)
card(sl,0.3,1.3,6.1,5.5,C_PURPLE); card(sl,6.9,1.3,6.1,5.5,C_RED)
p1=add_pic(sl,PICS[0],0.45,1.45,5.8,3.9)
p2=add_pic(sl,PICS[1],7.05,1.45,5.8,3.9)
txt(sl,"Home View — Recently Played & Recommendations",0.4,5.45,5.9,0.6,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
txt(sl,"Library View — Song Table with Artwork",7.0,5.45,5.9,0.6,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
if p1: add_anim(sl,p1,200)
if p2: add_anim(sl,p2,400)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Screenshots 3 & 4
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Now Playing & Liked Songs","Full-screen overlay and liked songs view"); footer_bar(sl,7)
transition_none(sl)
card(sl,0.3,1.3,6.1,5.5,C_CYAN); card(sl,6.9,1.3,6.1,5.5,C_PURPLE)
p3=add_pic(sl,PICS[2],0.45,1.45,5.8,3.9)
p4=add_pic(sl,PICS[3],7.05,1.45,5.8,3.9)
txt(sl,"Now Playing — Full-screen artwork overlay",0.4,5.45,5.9,0.6,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
txt(sl,"Liked Songs — Filtered heart-toggled songs",7.0,5.45,5.9,0.6,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
if p3: add_anim(sl,p3,200)
if p4: add_anim(sl,p4,400)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Playback
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Music Playback & Controls","PlayerController + PlayerBar"); footer_bar(sl,8)
transition_none(sl)
c1=card(sl,0.3,1.3,6.1,5.5,C_RED); c2=card(sl,6.8,1.3,6.2,5.5,C_PURPLE)
txt(sl,"PlayerController",0.5,1.38,5.7,0.45,size=15,bold=True,color=C_RED)
txt(sl,"PlayerBar Component",7.0,1.38,5.8,0.45,size=15,bold=True,color=C_PURPLE)
b1=bullets(sl,["JavaFX MediaPlayer for audio","Play / Pause / Next / Previous","Shuffle — random song selection","Repeat — loop current song","Seek bar — jump to position","Volume control slider","##Callbacks","onSongChange, onLikeChange","onPlayStateChange"],0.45,1.85,5.8,4.8)
b2=bullets(sl,["Always visible at bottom","Shows artwork thumbnail","Song title + artist name","Progress bar with time","Heart button for like/unlike","Volume slider","##Design","Listener interface pattern","Decoupled from controller"],6.95,1.85,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: add_anim(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Design Patterns
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Design Patterns Used","Clean architecture through proven patterns"); footer_bar(sl,9)
transition_none(sl)
patterns=[
    ("MVC","Model-View-Controller separates data, UI, and logic cleanly",C_RED),
    ("Observer","Callbacks: onSongChange, onLikeChange, onPlayStateChange",C_PURPLE),
    ("Strategy","Shuffle vs Sequential playback — swappable at runtime",C_CYAN),
    ("Facade","DatabaseService hides mp3agic + MediaPlayer complexity",C_RED),
    ("Factory","SongTableListener / MusicCardListener interface factories",C_PURPLE),
    ("Singleton","DatabaseService — single source of truth for all song data",C_CYAN),
]
for i,(name,desc,col) in enumerate(patterns):
    row=i//2; ci=i%2
    fx=0.3+ci*6.5; fy=1.35+row*1.85
    r=card(sl,fx,fy,6.2,1.65,col)
    rect(sl,fx,fy,0.15,1.65,col)
    txt(sl,name,fx+0.3,fy+0.15,2.5,0.55,size=18,bold=True,color=col,font="Calibri")
    txt(sl,desc,fx+0.3,fy+0.78,5.7,0.7,size=13,color=C_LGRAY)
    add_anim(sl,r,i*120)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Persistence
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Data Persistence & Storage","Library survives app restarts"); footer_bar(sl,10)
transition_none(sl)
r0=card(sl,0.3,1.3,12.73,1.05,C_PURPLE)
txt(sl,"Save File:  %APPDATA%\\AuralinPlayer\\library.dat",0.5,1.35,12.3,0.42,size=14,bold=True,color=C_WHITE)
txt(sl,"filePath  |  liked  |  plays  |  lastPlayed  |  artist  |  durationMs",0.5,1.78,12.3,0.42,size=13,color=C_GOLD,font="Courier New")
c1=card(sl,0.3,2.5,6.1,4.3,C_RED); c2=card(sl,6.8,2.5,6.2,4.3,C_PURPLE)
txt(sl,"Save Triggers",0.5,2.58,5.7,0.42,size=15,bold=True,color=C_RED)
txt(sl,"Load on Startup",7.0,2.58,5.8,0.42,size=15,bold=True,color=C_PURPLE)
b1=bullets(sl,["On song play — recordPlay()","On like toggle — saveToDisk()","On app close — stop()","Background thread — non-blocking","##Why Fast?","mp3agic reads tags directly","No MediaPlayer on startup","Thread pool (3 threads)"],0.45,3.05,5.8,3.6)
b2=bullets(sl,["Read pipe-separated lines","Skip missing files automatically","Restore: liked, plays, lastPlayed","Restore: artist, durationMs","##Artwork Loading","Background thread per song","Platform.runLater for UI update","Table refreshes incrementally"],6.95,3.05,5.9,3.6)
for s,d in [(r0,0),(c1,100),(c2,200),(b1,300),(b2,400)]: add_anim(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — UI/UX
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"UI/UX Design Highlights","Premium dark aesthetic inspired by Spotify"); footer_bar(sl,11)
transition_none(sl)
colors_data=[("#050510","BG"),("#101022","Card"),("#FA2D48","Red"),("#7B2FFF","Purple"),("#00D4FF","Cyan"),("#FFC800","Gold")]
for i,(hx,nm) in enumerate(colors_data):
    fx=0.3+i*2.15; fy=1.35
    r=rect(sl,fx,fy,1.9,1.0,RGBColor(int(hx[1:3],16),int(hx[3:5],16),int(hx[5:7],16)))
    txt(sl,nm,fx,fy+1.05,1.9,0.38,size=12,color=C_LGRAY,align=PP_ALIGN.CENTER)
    txt(sl,hx,fx,fy+1.42,1.9,0.32,size=10,color=C_MGRAY,align=PP_ALIGN.CENTER,italic=True)
c1=card(sl,0.3,2.95,6.1,3.85,C_RED); c2=card(sl,6.8,2.95,6.2,3.85,C_PURPLE)
txt(sl,"Custom Window",0.5,3.03,5.7,0.42,size=15,bold=True,color=C_RED)
txt(sl,"Components",7.0,3.03,5.8,0.42,size=15,bold=True,color=C_PURPLE)
bullets(sl,["StageStyle.TRANSPARENT","Custom title bar (32px)","Win11-style min/max/close","8-edge resize handles","Rounded corners (radius 10)","Screen-aware sizing on start"],0.45,3.5,5.8,3.1)
bullets(sl,["MusicCard — 160px hover card","SongTable — custom TableView","PlayerBar — bottom fixed bar","Sidebar — navigation panel","NowPlayingView — overlay","ScrollPane hidden scrollbars"],6.95,3.5,5.9,3.1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Challenges
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Challenges & Solutions","Real problems solved during development"); footer_bar(sl,12)
transition_none(sl)
chs=[
    ("Slow startup","MediaPlayer for each song on load","mp3agic reads tags; no MediaPlayer needed"),
    ("Scrollbar appearing","TableView internal ScrollPane","skinProperty listener hides bars post-skin"),
    ("Missing DLLs in EXE","jlink doesn't copy JavaFX native DLLs","Copy javafx/bin/*.dll into runtime/bin/"),
    ("Data not persisting","No save mechanism existed","Pipe-separated file with all metadata cached"),
    ("Open With missing","App not registered in Windows","Inno Setup registry + FriendlyAppName key"),
]
cols_alt=[C_RED,C_PURPLE,C_CYAN,C_RED,C_PURPLE]
for i,(prob,cause,sol) in enumerate(chs):
    fy=1.35+i*1.1
    r=card(sl,0.3,fy,12.73,0.98,cols_alt[i])
    rect(sl,0.3,fy,0.15,0.98,cols_alt[i])
    txt(sl,prob,0.55,fy+0.08,2.8,0.42,size=14,bold=True,color=cols_alt[i])
    txt(sl,"Cause: "+cause,3.5,fy+0.08,4.5,0.38,size=12,color=C_LGRAY,italic=True)
    txt(sl,"Fix: "+sol,8.2,fy+0.08,4.7,0.38,size=12,color=C_WHITE)
    add_anim(sl,r,i*150)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Installer
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Installer & Deployment","Zero-dependency Windows installer"); footer_bar(sl,13)
transition_none(sl)
steps=[("1","Compile Java","javac -> .class",C_RED),("2","Fat JAR","jar + mp3agic",C_PURPLE),("3","jlink JRE","80MB bundled JRE",C_CYAN),("4","launch4j","JAR -> .exe",C_RED),("5","Inno Setup","Full installer",C_PURPLE)]
for i,(num,title,desc,col) in enumerate(steps):
    fx=0.3+i*2.6; fy=1.4
    r=card(sl,fx,fy,2.4,2.2,col)
    rect(sl,fx,fy,2.4,0.5,col)
    txt(sl,num,fx+0.05,fy+0.06,0.45,0.38,size=18,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER)
    txt(sl,title,fx+0.1,fy+0.58,2.2,0.48,size=14,bold=True,color=C_WHITE)
    txt(sl,desc,fx+0.1,fy+1.12,2.2,0.85,size=12,color=C_LGRAY)
    if i<4: txt(sl,"->",fx+2.4,fy+0.85,0.2,0.4,size=18,bold=True,color=C_GOLD,align=PP_ALIGN.CENTER)
    add_anim(sl,r,i*200)
r2=card(sl,0.3,3.8,12.73,2.95,C_VIOLET)
txt(sl,"What the installer does",0.5,3.88,12.3,0.42,size=15,bold=True,color=C_WHITE)
bullets(sl,["Installs AuralinPlayer.exe + bundled runtime to Program Files","Registers .mp3 Open With association in Windows registry","Creates Desktop + Start Menu shortcuts with custom icon","Includes EULA acceptance screen  |  Full uninstall support"],0.5,4.38,12.3,2.2)
add_anim(sl,r2,1000)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl); header_bar(sl,"Conclusion & Future Work","What was achieved and what comes next"); footer_bar(sl,14)
transition_none(sl)
c1=card(sl,0.3,1.3,6.1,5.5,C_RED); c2=card(sl,6.8,1.3,6.2,5.5,C_PURPLE)
txt(sl,"Achievements",0.5,1.38,5.7,0.45,size=15,bold=True,color=C_RED)
txt(sl,"Future Work",7.0,1.38,5.8,0.45,size=15,bold=True,color=C_PURPLE)
b1=bullets(sl,["Full-featured music player","Modern Spotify-inspired dark UI","MVC architecture with patterns","Persistent library — fast startup","Bundled JRE installer","Open With MP3 registration","Custom frameless Win11 window","Background metadata loading"],0.45,1.85,5.8,4.8)
b2=bullets(sl,["Equalizer & audio effects","Online streaming integration","Playlist & queue management","Cross-platform Linux/macOS","Album & artist grouping","Lyrics display integration","Mini player mode","Cloud sync for library"],6.95,1.85,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: add_anim(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
premium_bg(sl)
rect(sl,0,0,13.33,7.5,C_VIOLET)
rect(sl,0,0,13.33,7.5,C_BG)
rect(sl,3.5,0.8,6.33,5.9,C_CARD2)
rect(sl,3.5,0.8,6.33,0.08,C_PURPLE)
rect(sl,3.5,6.62,6.33,0.08,C_RED)
rect(sl,3.5,0.8,0.08,5.9,C_PURPLE)
rect(sl,9.75,0.8,0.08,5.9,C_RED)
t1=txt(sl,"♪",5.5,1.2,2.5,2.5,size=88,bold=True,color=C_PURPLE,align=PP_ALIGN.CENTER)
t2=txt(sl,"Thank You!",3.8,3.8,5.8,1.1,size=50,bold=True,color=C_WHITE,align=PP_ALIGN.CENTER,font="Calibri")
t3=txt(sl,"Auralin Music Player",4.0,4.95,5.33,0.55,size=18,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
t4=txt(sl,"Dip Karmokar   |   OOP Project   |   2026",3.8,5.55,5.73,0.45,size=14,color=C_LGRAY,align=PP_ALIGN.CENTER)
rect(sl,0,6.75,13.33,0.75,C_CARD)
rect(sl,0,6.75,13.33,0.06,C_PURPLE)
txt(sl,"Questions? Let's discuss!",0.5,6.85,12.33,0.45,size=14,color=C_LGRAY,align=PP_ALIGN.CENTER,italic=True)
transition_none(sl)
for s,d in [(t1,0),(t2,300),(t3,500),(t4,700)]: add_anim(sl,s,d)

# ── Save ──────────────────────────────────────────────────────────────────
out = r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\AuralinPresentation.pptx"
prs.save(out)
print("Saved:", out)
