from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Premium Dark Theme Colors ─────────────────────────────────────────────
BG1    = RGBColor(0x0A, 0x0A, 0x0F)   # near black
BG2    = RGBColor(0x12, 0x12, 0x20)   # dark navy
CARD   = RGBColor(0x1E, 0x1E, 0x35)   # card bg
GRAD1  = RGBColor(0x6C, 0x00, 0xC8)   # deep violet
GRAD2  = RGBColor(0xFA, 0x2D, 0x48)   # Auralin red
ACCENT = RGBColor(0xA0, 0x40, 0xFF)   # bright purple
GOLD   = RGBColor(0xFF, 0xD7, 0x00)   # gold highlight
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xB0, 0xB0, 0xCC)
PINK   = RGBColor(0xFF, 0x80, 0xAB)

PICS = [
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (1).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (2).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (3).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (4).png",
]

# ── Animation helpers ─────────────────────────────────────────────────────
def add_slide_transition(slide, effect="fade"):
    """Add morph/fade transition to slide."""
    spTree = slide.shapes._spTree
    mc = etree.SubElement(slide._element, qn('p:transition'))
    mc.set('spd', 'med')
    mc.set('advTm', '0')
    if effect == "fade":
        etree.SubElement(mc, qn('p:fade'))
    elif effect == "push":
        push = etree.SubElement(mc, qn('p:push'))
        push.set('dir', 'l')
    elif effect == "wipe":
        wipe = etree.SubElement(mc, qn('p:wipe'))
        wipe.set('dir', 'r')

def make_anim_seq(slide):
    """Return (or create) the timing/tnLst element for animations."""
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(slide._element, qn('p:timing'))
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = etree.SubElement(timing, qn('p:tnLst'))
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = etree.SubElement(tnLst, qn('p:par'))
        cTn = etree.SubElement(par, qn('p:cTn'))
        cTn.set('id','1'); cTn.set('dur','indefinite'); cTn.set('restart','whenNotActive')
        cTn.set('nodeType','tmRoot')
        childTnLst = etree.SubElement(cTn, qn('p:childTnLst'))
        seq = etree.SubElement(childTnLst, qn('p:seq'))
        seq.set('concurrent','1'); seq.set('nextAc','seek')
        cTn2 = etree.SubElement(seq, qn('p:cTn'))
        cTn2.set('id','2'); cTn2.set('dur','indefinite'); cTn2.set('nodeType','mainSeq')
        etree.SubElement(cTn2, qn('p:childTnLst'))
        etree.SubElement(seq, qn('p:prevCondLst'))
        etree.SubElement(seq, qn('p:nextCondLst'))
    return par.find(qn('p:cTn')).find(qn('p:childTnLst')).find(qn('p:seq')).find(qn('p:cTn')).find(qn('p:childTnLst'))

def add_appear_anim(slide, shape, delay_ms=0, effect="fade"):
    """Add fade/fly-in animation to shape."""
    try:
        childTnLst = make_anim_seq(slide)
        idx = len(childTnLst) + 3
        par = etree.SubElement(childTnLst, qn('p:par'))
        cTn = etree.SubElement(par, qn('p:cTn'))
        cTn.set('id', str(idx)); cTn.set('fill','hold')
        stCondLst = etree.SubElement(cTn, qn('p:stCondLst'))
        cond = etree.SubElement(stCondLst, qn('p:cond'))
        cond.set('delay', str(delay_ms))
        childTnLst2 = etree.SubElement(cTn, qn('p:childTnLst'))
        par2 = etree.SubElement(childTnLst2, qn('p:par'))
        cTn2 = etree.SubElement(par2, qn('p:cTn'))
        cTn2.set('id', str(idx+1)); cTn2.set('fill','hold')
        stCondLst2 = etree.SubElement(cTn2, qn('p:stCondLst'))
        cond2 = etree.SubElement(stCondLst2, qn('p:cond'))
        cond2.set('delay','0')
        childTnLst3 = etree.SubElement(cTn2, qn('p:childTnLst'))
        animEffect = etree.SubElement(childTnLst3, qn('p:animEffect'))
        animEffect.set('transition','in'); animEffect.set('filter','fade')
        cBhvr = etree.SubElement(animEffect, qn('p:cBhvr'))
        cTn3 = etree.SubElement(cBhvr, qn('p:cTn'))
        cTn3.set('id', str(idx+2)); cTn3.set('dur','500'); cTn3.set('fill','hold')
        tgtEl = etree.SubElement(cBhvr, qn('p:tgtEl'))
        spTgt = etree.SubElement(tgtEl, qn('p:spTgt'))
        spTgt.set('spid', str(shape.shape_id))
    except Exception:
        pass

# ── Drawing helpers ───────────────────────────────────────────────────────
def bg(slide, color=BG1):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = font
    return tb

def bullets(slide, items, l, t, w, h, base_size=17):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if item.startswith("##"):
            r = p.add_run(); r.text = item[2:]
            r.font.size = Pt(base_size+3); r.font.bold = True
            r.font.color.rgb = PINK; r.font.name = "Montserrat"
        else:
            lvl = 1 if item.startswith("  ") else 0
            p.level = lvl
            r = p.add_run()
            r.text = ("▸ " if lvl==0 else "   · ") + item.strip()
            r.font.size = Pt(base_size if lvl==0 else base_size-1)
            r.font.color.rgb = WHITE if lvl==0 else LGRAY
            r.font.name = "Calibri"
    return tb

def line(slide, y, color=ACCENT, thickness=0.04):
    s = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Inches(thickness))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.15, CARD)
    rect(slide, 0, 0, 0.18, 1.15, GRAD2)
    t = txt(slide, title, 0.4, 0.1, 11, 0.9, size=34, bold=True, color=WHITE, font="Montserrat")
    if subtitle:
        txt(slide, subtitle, 0.4, 0.88, 11, 0.35, size=13, color=LGRAY, italic=True)
    line(slide, 1.18, GRAD1, 0.05)
    txt(slide, "♪", 12.3, 0.1, 0.9, 0.9, size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return t

def footer(slide, page_num):
    rect(slide, 0, 7.1, 13.33, 0.4, CARD)
    rect(slide, 0, 7.1, 4, 0.4, GRAD2)
    txt(slide, "Auralin Music Player  |  Dip Karmokar", 0.2, 7.12, 5, 0.3, size=11, color=WHITE)
    txt(slide, f"{page_num} / 15", 12.5, 7.12, 0.8, 0.3, size=11, color=LGRAY, align=PP_ALIGN.RIGHT)

def add_pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        # rounded border
        sp = pic._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(sp, qn('p:spPr'))
        ln = etree.SubElement(spPr, qn('a:ln'))
        ln.set('w','19050')
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', 'A040FF')
        return pic
    return None

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
# Gradient blobs
rect(sl, -1, -1, 7, 5, GRAD1); rect(sl, 7, 3, 7, 5, GRAD2)
rect(sl, 1, 1, 11.33, 5.5, BG1)  # dark overlay
# Glowing circle
rect(sl, 9.5, 0.5, 3.5, 3.5, ACCENT)
rect(sl, 9.8, 0.8, 2.9, 2.9, BG1)
t1 = txt(sl, "♪", 10.0, 0.9, 2.5, 2.5, size=90, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
# Title
t2 = txt(sl, "AURALIN", 0.8, 1.2, 9, 1.5, size=80, bold=True, color=WHITE, font="Montserrat")
t3 = txt(sl, "MUSIC PLAYER", 0.8, 2.6, 9, 1.0, size=44, bold=True, color=GRAD2, font="Montserrat")
line(sl, 3.75, ACCENT, 0.06)
t4 = txt(sl, "A Modern JavaFX Desktop Music Application", 0.8, 3.9, 10, 0.6, size=20, color=LGRAY, italic=True)
t5 = txt(sl, "Dip Karmokar   ·   OOP Project   ·   2026", 0.8, 4.6, 10, 0.5, size=16, color=LGRAY)
rect(sl, 0, 6.8, 13.33, 0.7, GRAD2)
txt(sl, "JavaFX 25  ·  mp3agic  ·  MVC Architecture  ·  Windows Installer", 1, 6.85, 11, 0.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_slide_transition(sl, "fade")
for s in [t2,t3,t4,t5]: add_appear_anim(sl, s, delay_ms=300)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
h = header(sl, "Agenda", "What we'll cover today")
footer(sl, 2)
add_slide_transition(sl, "push")
agenda = [
    ("01", "Introduction & Motivation",    GRAD2),
    ("02", "System Architecture",          ACCENT),
    ("03", "Key Features",                 GRAD1),
    ("04", "Implementation & Patterns",    GRAD2),
    ("05", "Data Persistence",             ACCENT),
    ("06", "UI/UX Design",                 GRAD1),
    ("07", "Challenges & Solutions",       GRAD2),
    ("08", "Conclusion & Future Work",     ACCENT),
]
cols = [(0.4, 1.3), (6.9, 1.3), (0.4, 3.8), (6.9, 3.8),
        (0.4, 5.5), (6.9, 5.5), (0.4, 6.5), (6.9, 6.5)]
for i, (num, label, col) in enumerate(agenda):
    cx, cy = cols[i]
    r = rect(sl, cx, cy, 6.0, 0.9, CARD)
    rect(sl, cx, cy, 0.55, 0.9, col)
    t = txt(sl, num, cx+0.05, cy+0.15, 0.5, 0.6, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t2 = txt(sl, label, cx+0.65, cy+0.18, 5.2, 0.55, size=17, bold=True, color=WHITE)
    add_appear_anim(sl, r, delay_ms=i*150)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Introduction", "What is Auralin Music Player?")
footer(sl, 3); add_slide_transition(sl, "fade")
rect(sl, 0.4, 1.3, 5.8, 5.5, CARD)
rect(sl, 7.0, 1.3, 5.9, 5.5, CARD)
rect(sl, 0.4, 1.3, 5.8, 0.45, GRAD2)
rect(sl, 7.0, 1.3, 5.9, 0.45, ACCENT)
txt(sl, "What is Auralin?", 0.6, 1.33, 5.4, 0.4, size=16, bold=True, color=WHITE)
txt(sl, "Goals & Scope", 7.2, 1.33, 5.5, 0.4, size=16, bold=True, color=WHITE)
b1 = bullets(sl, [
    "Modern desktop music player",
    "Built with JavaFX 25 + Java",
    "Inspired by Spotify dark UI",
    "##Motivation",
    "Existing players lack modern UI",
    "Deep dive into JavaFX MVC",
    "##Tech Stack",
    "JavaFX 25, mp3agic, Inno Setup",
    "launch4j, jlink bundled JRE",
], 0.55, 1.85, 5.5, 4.8)
b2 = bullets(sl, [
    "Fast startup — no MediaPlayer on load",
    "Persistent library across sessions",
    "Open With — default MP3 player",
    "Custom frameless Win11 window",
    "##Deliverables",
    "Runnable .exe installer",
    "Bundled JRE — no Java needed",
    "Full MVC architecture",
], 7.15, 1.85, 5.6, 4.8)
for s in [b1, b2]: add_appear_anim(sl, s, delay_ms=200)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "System Architecture", "MVC Pattern with Observer & Facade")
footer(sl, 4); add_slide_transition(sl, "wipe")
arch_boxes = [
    (0.3,  1.4, 3.8, 5.3, GRAD2,  "MODEL",       ["Song.java","title, artist","filePath","liked, plays","lastPlayed","getDurationStr()"]),
    (4.75, 1.4, 3.8, 5.3, ACCENT, "VIEW",         ["HomeView","LibraryView","SearchView","LikedView","SongTable","PlayerBar","Sidebar","NowPlayingView"]),
    (9.2,  1.4, 3.8, 5.3, GRAD1,  "CONTROLLER",  ["PlayerController","ViewManager","DatabaseService","FileImportService","PlayerBar.Listener","MusicCard.Listener"]),
]
for bx, by, bw, bh, bc, bt, bi in arch_boxes:
    r = rect(sl, bx, by, bw, bh, CARD)
    rect(sl, bx, by, bw, 0.5, bc)
    txt(sl, bt, bx+0.1, by+0.08, bw-0.2, 0.38, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Montserrat")
    for j, item in enumerate(bi):
        txt(sl, "· "+item, bx+0.2, by+0.6+j*0.7, bw-0.3, 0.6, size=13, color=LGRAY)
    add_appear_anim(sl, r, delay_ms=200)
txt(sl, "⟷", 4.1, 3.8, 0.6, 0.5, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "⟷", 8.6, 3.8, 0.6, 0.5, size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "Observer  ·  MVC  ·  Facade  ·  Strategy  ·  Singleton", 2, 6.85, 9.33, 0.4, size=13, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Key Features", "Everything Auralin can do")
footer(sl, 5); add_slide_transition(sl, "fade")
feats = [
    ("🎵", "Library",      "Add files & folders",   GRAD2),
    ("🔍", "Search",       "Real-time filtering",    ACCENT),
    ("❤️",  "Liked Songs",  "Heart toggle + persist", GRAD1),
    ("🏠", "Home View",    "Recent + Recommended",   GRAD2),
    ("🎬", "Now Playing",  "Full-screen overlay",    ACCENT),
    ("💾", "Auto-Save",    "Persists across sessions",GRAD1),
    ("📂", "Open With",    "Default MP3 player",     GRAD2),
    ("🖥️",  "Custom Window","Frameless + Win11 btns", ACCENT),
]
positions = [(0.3,1.3),(3.55,1.3),(6.8,1.3),(10.05,1.3),
             (0.3,3.9),(3.55,3.9),(6.8,3.9),(10.05,3.9)]
for i, ((fx,fy),(icon,title,desc,col)) in enumerate(zip(positions,feats)):
    r = rect(sl, fx, fy, 3.0, 2.3, CARD)
    rect(sl, fx, fy, 3.0, 0.08, col)
    txt(sl, icon, fx+0.1, fy+0.15, 0.7, 0.7, size=28)
    txt(sl, title, fx+0.85, fy+0.18, 2.0, 0.5, size=15, bold=True, color=WHITE, font="Montserrat")
    txt(sl, desc, fx+0.1, fy+0.85, 2.8, 0.9, size=13, color=LGRAY)
    add_appear_anim(sl, r, delay_ms=i*100)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Screenshot: Home View
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Home View & Library", "Screenshots from Auralin Music Player")
footer(sl, 6); add_slide_transition(sl, "fade")
rect(sl, 0.3, 1.3, 6.1, 5.5, CARD)
rect(sl, 6.9, 1.3, 6.1, 5.5, CARD)
p1 = add_pic(sl, PICS[0], 0.4, 1.4, 5.9, 4.0)
p2 = add_pic(sl, PICS[1], 7.0, 1.4, 5.9, 4.0)
txt(sl, "Home View — Recently Played & Recommendations", 0.4, 5.5, 5.9, 0.6, size=13, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
txt(sl, "Library View — Song Table with Artwork", 7.0, 5.5, 5.9, 0.6, size=13, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
if p1: add_appear_anim(sl, p1, delay_ms=200)
if p2: add_appear_anim(sl, p2, delay_ms=400)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Screenshot: Now Playing & Liked
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Now Playing & Liked Songs", "Full-screen overlay and liked songs view")
footer(sl, 7); add_slide_transition(sl, "fade")
rect(sl, 0.3, 1.3, 6.1, 5.5, CARD)
rect(sl, 6.9, 1.3, 6.1, 5.5, CARD)
p3 = add_pic(sl, PICS[2], 0.4, 1.4, 5.9, 4.0)
p4 = add_pic(sl, PICS[3], 7.0, 1.4, 5.9, 4.0)
txt(sl, "Now Playing — Full-screen artwork overlay", 0.4, 5.5, 5.9, 0.6, size=13, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
txt(sl, "Liked Songs — Filtered heart-toggled songs", 7.0, 5.5, 5.9, 0.6, size=13, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
if p3: add_appear_anim(sl, p3, delay_ms=200)
if p4: add_appear_anim(sl, p4, delay_ms=400)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Playback & Controls
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Music Playback & Controls", "PlayerController + PlayerBar")
footer(sl, 8); add_slide_transition(sl, "push")
rect(sl, 0.3, 1.3, 6.0, 5.5, CARD); rect(sl, 6.8, 1.3, 6.2, 5.5, CARD)
rect(sl, 0.3, 1.3, 6.0, 0.45, GRAD2); rect(sl, 6.8, 1.3, 6.2, 0.45, ACCENT)
txt(sl, "PlayerController", 0.5, 1.33, 5.6, 0.4, size=15, bold=True, color=WHITE)
txt(sl, "PlayerBar Component", 7.0, 1.33, 5.8, 0.4, size=15, bold=True, color=WHITE)
b1 = bullets(sl, ["JavaFX MediaPlayer for audio","Play / Pause / Next / Previous","Shuffle — random song selection","Repeat — loop current song","Seek bar — jump to position","Volume control slider","##Callbacks","onSongChange, onLikeChange","onPlayStateChange"], 0.45, 1.85, 5.7, 4.8)
b2 = bullets(sl, ["Always visible at bottom","Shows artwork thumbnail","Song title + artist name","Progress bar with time","Heart button for like/unlike","Volume slider","##Design","Listener interface pattern","Decoupled from controller"], 6.95, 1.85, 5.9, 4.8)
for s in [b1,b2]: add_appear_anim(sl, s, delay_ms=200)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Implementation & Design Patterns
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Design Patterns Used", "Clean architecture through proven patterns")
footer(sl, 9); add_slide_transition(sl, "fade")
patterns = [
    ("MVC",      "Model-View-Controller separates data, UI, and logic cleanly", GRAD2),
    ("Observer", "Callbacks: onSongChange, onLikeChange, onPlayStateChange",    ACCENT),
    ("Strategy", "Shuffle vs Sequential playback — swappable at runtime",       GRAD1),
    ("Facade",   "DatabaseService hides mp3agic + MediaPlayer complexity",      GRAD2),
    ("Factory",  "SongTableListener / MusicCardListener interface factories",   ACCENT),
    ("Singleton","DatabaseService — single source of truth for all song data",  GRAD1),
]
for i, (name, desc, col) in enumerate(patterns):
    row = i // 2; col_idx = i % 2
    fx = 0.3 + col_idx * 6.5; fy = 1.35 + row * 1.85
    r = rect(sl, fx, fy, 6.2, 1.6, CARD)
    rect(sl, fx, fy, 0.12, 1.6, col)
    txt(sl, name, fx+0.25, fy+0.12, 2.5, 0.55, size=18, bold=True, color=col, font="Montserrat")
    txt(sl, desc, fx+0.25, fy+0.72, 5.8, 0.75, size=13, color=LGRAY)
    add_appear_anim(sl, r, delay_ms=i*120)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Persistence
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Data Persistence & Storage", "Library survives app restarts")
footer(sl, 10); add_slide_transition(sl, "wipe")
rect(sl, 0.3, 1.3, 12.73, 1.1, CARD)
rect(sl, 0.3, 1.3, 12.73, 0.45, GRAD1)
txt(sl, "Save File Format  —  %APPDATA%\\AuralinPlayer\\library.dat", 0.5, 1.33, 12.3, 0.4, size=15, bold=True, color=WHITE)
txt(sl, "filePath  |  liked  |  plays  |  lastPlayed  |  artist  |  durationMs", 0.5, 1.82, 12.3, 0.45, size=14, color=GOLD, font="Courier New")
rect(sl, 0.3, 2.55, 6.0, 4.2, CARD); rect(sl, 6.8, 2.55, 6.2, 4.2, CARD)
rect(sl, 0.3, 2.55, 6.0, 0.45, GRAD2); rect(sl, 6.8, 2.55, 6.2, 0.45, ACCENT)
txt(sl, "Save Triggers", 0.5, 2.58, 5.6, 0.4, size=15, bold=True, color=WHITE)
txt(sl, "Load on Startup", 7.0, 2.58, 5.8, 0.4, size=15, bold=True, color=WHITE)
b1 = bullets(sl, ["On song play — recordPlay()","On like toggle — saveToDisk()","On app close — stop()","Background thread — non-blocking","##Why Fast?","mp3agic reads tags directly","No MediaPlayer on startup","Thread pool (3 threads)"], 0.45, 3.1, 5.7, 3.5)
b2 = bullets(sl, ["Read pipe-separated lines","Skip missing files automatically","Restore: liked, plays, lastPlayed","Restore: artist, durationMs","##Artwork Loading","Background thread per song","Platform.runLater for UI update","Table refreshes incrementally"], 6.95, 3.1, 5.9, 3.5)
for s in [b1,b2]: add_appear_anim(sl, s, delay_ms=200)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — UI/UX Design
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "UI/UX Design Highlights", "Premium dark aesthetic inspired by Spotify")
footer(sl, 11); add_slide_transition(sl, "fade")
colors_data = [("#0A0A0F","Background"),("#1E1E35","Card BG"),("#FA2D48","Auralin Red"),("#A040FF","Accent Purple"),("#B3B3B3","Text Gray"),("#FFD700","Highlight")]
for i, (hex_c, name) in enumerate(colors_data):
    fx = 0.3 + i*2.15; fy = 1.35
    r = rect(sl, fx, fy, 1.9, 1.0, RGBColor(int(hex_c[1:3],16),int(hex_c[3:5],16),int(hex_c[5:7],16)))
    txt(sl, name, fx, fy+1.05, 1.9, 0.4, size=11, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(sl, hex_c, fx, fy+1.45, 1.9, 0.35, size=10, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)
rect(sl, 0.3, 2.95, 6.0, 3.8, CARD); rect(sl, 6.8, 2.95, 6.2, 3.8, CARD)
rect(sl, 0.3, 2.95, 6.0, 0.45, GRAD2); rect(sl, 6.8, 2.95, 6.2, 0.45, ACCENT)
txt(sl, "Custom Window", 0.5, 2.98, 5.6, 0.4, size=15, bold=True, color=WHITE)
txt(sl, "Components", 7.0, 2.98, 5.8, 0.4, size=15, bold=True, color=WHITE)
bullets(sl, ["StageStyle.TRANSPARENT","Custom title bar (32px)","Win11-style min/max/close","8-edge resize handles","Rounded corners (radius 10)","Screen-aware sizing on start"], 0.45, 3.5, 5.7, 3.1)
bullets(sl, ["MusicCard — 160px hover card","SongTable — custom TableView","PlayerBar — bottom fixed bar","Sidebar — navigation panel","NowPlayingView — overlay","ScrollPane hidden scrollbars"], 6.95, 3.5, 5.9, 3.1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Challenges & Solutions
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Challenges & Solutions", "Real problems solved during development")
footer(sl, 12); add_slide_transition(sl, "push")
challenges = [
    ("Slow startup",        "MediaPlayer for each song on load",    "mp3agic reads tags; no MediaPlayer needed"),
    ("Scrollbar appearing", "TableView internal ScrollPane",        "skinProperty listener hides bars post-skin"),
    ("Missing DLLs in EXE", "jlink doesn't copy JavaFX native DLLs","Copy javafx/bin/*.dll into runtime/bin/"),
    ("Data not persisting", "No save mechanism existed",            "Pipe-separated file with all metadata cached"),
    ("Open With missing",   "App not registered in Windows",        "Inno Setup registry + FriendlyAppName key"),
]
for i, (prob, cause, sol) in enumerate(challenges):
    fy = 1.35 + i * 1.1
    r = rect(sl, 0.3, fy, 12.73, 0.95, CARD)
    rect(sl, 0.3, fy, 0.12, 0.95, GRAD2 if i%2==0 else ACCENT)
    txt(sl, "⚠ "+prob, 0.55, fy+0.05, 2.8, 0.4, size=14, bold=True, color=GRAD2 if i%2==0 else ACCENT)
    txt(sl, "Cause: "+cause, 3.5, fy+0.05, 4.5, 0.38, size=12, color=LGRAY, italic=True)
    txt(sl, "✓ "+sol, 8.2, fy+0.05, 4.7, 0.38, size=12, color=WHITE)
    add_appear_anim(sl, r, delay_ms=i*150)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Installer & Deployment
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Installer & Deployment", "Zero-dependency Windows installer")
footer(sl, 13); add_slide_transition(sl, "fade")
steps = [("1","Compile Java","javac → .class files",GRAD2),("2","Build Fat JAR","jar + mp3agic extracted",ACCENT),("3","jlink JRE","Custom 80MB bundled JRE",GRAD1),("4","launch4j EXE","JAR wrapped as .exe",GRAD2),("5","Inno Setup","Full installer with registry",ACCENT)]
for i, (num, title, desc, col) in enumerate(steps):
    fx = 0.3 + i*2.6; fy = 1.4
    r = rect(sl, fx, fy, 2.4, 2.2, CARD)
    rect(sl, fx, fy, 2.4, 0.5, col)
    txt(sl, num, fx+0.05, fy+0.05, 0.45, 0.4, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, fx+0.1, fy+0.55, 2.2, 0.5, size=14, bold=True, color=WHITE)
    txt(sl, desc, fx+0.1, fy+1.1, 2.2, 0.9, size=12, color=LGRAY)
    if i < 4:
        txt(sl, "→", fx+2.4, fy+0.85, 0.2, 0.4, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_appear_anim(sl, r, delay_ms=i*200)
rect(sl, 0.3, 3.8, 12.73, 2.9, CARD)
rect(sl, 0.3, 3.8, 12.73, 0.45, GRAD1)
txt(sl, "What the installer does", 0.5, 3.83, 12.3, 0.4, size=15, bold=True, color=WHITE)
bullets(sl, ["Installs AuralinPlayer.exe + bundled runtime to Program Files","Registers .mp3 Open With association in Windows registry","Creates Desktop + Start Menu shortcuts with custom icon","Includes EULA acceptance screen","Full uninstall support"], 0.5, 4.35, 12.3, 2.2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
header(sl, "Conclusion & Future Work", "What was achieved and what comes next")
footer(sl, 14); add_slide_transition(sl, "fade")
rect(sl, 0.3, 1.3, 6.0, 5.5, CARD); rect(sl, 6.8, 1.3, 6.2, 5.5, CARD)
rect(sl, 0.3, 1.3, 6.0, 0.45, GRAD2); rect(sl, 6.8, 1.3, 6.2, 0.45, ACCENT)
txt(sl, "✅  Achievements", 0.5, 1.33, 5.6, 0.4, size=15, bold=True, color=WHITE)
txt(sl, "🚀  Future Work", 7.0, 1.33, 5.8, 0.4, size=15, bold=True, color=WHITE)
b1 = bullets(sl, ["Full-featured music player","Modern Spotify-inspired dark UI","MVC architecture with patterns","Persistent library — fast startup","Bundled JRE installer","Open With MP3 registration","Custom frameless Win11 window","Background metadata loading"], 0.45, 1.85, 5.7, 4.8)
b2 = bullets(sl, ["Equalizer & audio effects","Online streaming integration","Playlist & queue management","Cross-platform Linux/macOS","Album & artist grouping","Lyrics display integration","Mini player mode","Cloud sync for library"], 6.95, 1.85, 5.9, 4.8)
for s in [b1,b2]: add_appear_anim(sl, s, delay_ms=200)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, BG1)
rect(sl, 0, 0, 13.33, 7.5, BG2)
rect(sl, 0, 0, 13.33, 7.5, GRAD1)
rect(sl, 1, 1, 11.33, 5.5, BG1)
rect(sl, 0, 6.5, 13.33, 1.0, GRAD2)
r1 = rect(sl, 4.67, 1.3, 4.0, 4.0, ACCENT)
r2 = rect(sl, 5.0, 1.6, 3.33, 3.4, BG1)
t1 = txt(sl, "♪", 5.5, 1.8, 2.5, 2.5, size=90, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
t2 = txt(sl, "Thank You!", 2.5, 5.4, 8.33, 1.0, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Montserrat")
t3 = txt(sl, "Auralin Music Player  —  A Modern JavaFX Desktop Music Application", 1.5, 6.1, 10.33, 0.5, size=14, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
t4 = txt(sl, "Dip Karmokar   ·   OOP Project   ·   2026", 3.5, 6.6, 6.33, 0.5, size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_slide_transition(sl, "fade")
for s in [t1,t2,t3,t4]: add_appear_anim(sl, s, delay_ms=300)

# ── Save ──────────────────────────────────────────────────────────────────
out = r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\AuralinPresentation.pptx"
prs.save(out)
print("✓ Saved:", out)
