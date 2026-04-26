from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette (matching purple/white music theme) ─────────────────────
BG_DARK    = RGBColor(0x1A, 0x00, 0x3D)   # deep purple
BG_MID     = RGBColor(0x3D, 0x00, 0x7A)   # medium purple
ACCENT     = RGBColor(0xC8, 0x00, 0xFF)   # bright purple/violet
ACCENT2    = RGBColor(0xFF, 0x2D, 0x48)   # red accent (Auralin brand)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xFF)   # light lavender

def add_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Montserrat" if bold else "Calibri"
    return txBox

def add_bullets(slide, items, l, t, w, h, size=18, title_color=ACCENT, body_color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if item.startswith("##"):
            run = p.add_run()
            run.text = item[2:]
            run.font.size = Pt(size + 2)
            run.font.bold = True
            run.font.color.rgb = title_color
            run.font.name = "Montserrat"
        else:
            p.level = 1 if item.startswith("  ") else 0
            run = p.add_run()
            run.text = ("• " if p.level == 0 else "  – ") + item.strip()
            run.font.size = Pt(size)
            run.font.color.rgb = body_color
            run.font.name = "Calibri"
    return txBox

def add_fade_animation(shape):
    """Add fade-in animation to a shape."""
    sp = shape._element
    timing = etree.SubElement(sp.getparent().getparent(), 
        '{http://schemas.openxmlformats.org/presentationml/2006/main}timing')

def divider_line(slide, y):
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

slides_data = [
    # Slide 1 — Title
    {"type": "title"},
    # Slide 2 — Agenda
    {"type": "content", "title": "Agenda",
     "items": ["Introduction & Motivation", "System Architecture", "Key Features",
               "Implementation Details", "Data Persistence", "UI/UX Design",
               "Challenges & Solutions", "Conclusion & Future Work"]},
    # Slide 3 — Introduction
    {"type": "content", "title": "Introduction",
     "items": ["##What is Auralin?",
               "A modern desktop music player built with JavaFX 25",
               "Inspired by Spotify's clean dark UI design",
               "##Motivation",
               "Existing players lack modern UI on desktop",
               "Learning JavaFX MVC architecture in depth",
               "##Goals",
               "Fast startup, smooth playback, persistent library",
               "Open With support — double-click MP3 to play"]},
    # Slide 4 — Architecture
    {"type": "arch"},
    # Slide 5 — Key Features
    {"type": "content", "title": "Key Features",
     "items": ["🎵  Library Management — add files & folders",
               "🔍  Search — real-time song filtering",
               "❤️   Liked Songs — heart toggle with persistence",
               "🏠  Home View — Recently Played & Recommendations",
               "🎬  Now Playing — full-screen artwork overlay",
               "💾  Auto-save — library persists across sessions",
               "📂  Open With — register as default MP3 player",
               "🖥️   Custom Window — frameless with Win11 controls"]},
    # Slide 6 — Library & File Import
    {"type": "content", "title": "Library Management & File Import",
     "items": ["##FileImportService",
               "Import individual MP3 files via file chooser",
               "Import entire folders recursively",
               "Progress bar shown during bulk import",
               "##DatabaseService",
               "mp3agic reads ID3v2 tags: artist, artwork, duration",
               "Background thread pool (3 threads) for fast loading",
               "No MediaPlayer needed on startup — instant open"]},
    # Slide 7 — Playback
    {"type": "content", "title": "Music Playback & Controls",
     "items": ["##PlayerController",
               "JavaFX MediaPlayer for audio playback",
               "Play / Pause / Next / Previous",
               "Shuffle mode — random song selection",
               "Repeat mode — repeat current song",
               "Seek bar — click to jump to position",
               "Volume control slider",
               "##PlayerBar Component",
               "Always visible at bottom of window",
               "Shows artwork, title, artist, progress"]},
    # Slide 8 — Home View
    {"type": "content", "title": "Home View & Recommendations",
     "items": ["##HomeView",
               "Greeting with username",
               "Recently Played — last 9 played songs",
               "Made For You — top 12 most played songs",
               "##MusicCard Component",
               "160px card with album artwork",
               "Hover scale animation (1.0 → 1.05)",
               "Red play/pause overlay for current song",
               "FlowPane grid — auto-reflows on resize"]},
    # Slide 9 — Liked & Now Playing
    {"type": "content", "title": "Liked Songs & Now Playing",
     "items": ["##Liked Songs View",
               "Filtered list of all liked songs",
               "Heart toggle in SongTable — instant update",
               "Unlike removes song from view immediately",
               "##Now Playing View",
               "Full overlay with large album artwork",
               "Blurred background effect",
               "Heart button to like/unlike current song",
               "Triggered by clicking artwork in PlayerBar"]},
    # Slide 10 — Design Patterns
    {"type": "content", "title": "Implementation & Design Patterns",
     "items": ["##MVC Architecture",
               "Model: Song.java — data only",
               "View: HomeView, LibraryView, SongTable...",
               "Controller: PlayerController, ViewManager",
               "##Observer Pattern",
               "onSongChange, onLikeChange, onPlayStateChange callbacks",
               "##Strategy Pattern",
               "Shuffle vs Sequential playback strategy",
               "##Facade Pattern",
               "DatabaseService hides mp3agic + MediaPlayer complexity"]},
    # Slide 11 — Persistence
    {"type": "content", "title": "Data Persistence & Storage",
     "items": ["##Save Location",
               "%APPDATA%\\AuralinPlayer\\library.dat",
               "##File Format (pipe-separated)",
               "filePath | liked | plays | lastPlayed | artist | durationMs",
               "##Save Triggers",
               "On song play (recordPlay)",
               "On like toggle (toggleLike / onLikeToggle)",
               "On app close (stop())",
               "##Load on Startup",
               "Instant — no MediaPlayer, reads from file",
               "Artwork loaded in background thread pool"]},
    # Slide 12 — UI Design
    {"type": "content", "title": "UI/UX Design Highlights",
     "items": ["##Color Palette",
               "Background: #111111 (near black)",
               "Accent: #FA2D48 (Auralin red)",
               "Text: #FFFFFF / #B3B3B3 (gray)",
               "##Custom Window (Frameless)",
               "StageStyle.TRANSPARENT — no OS chrome",
               "Custom title bar with Win11-style buttons",
               "8-edge resize handles",
               "##Responsive Layout",
               "Screen-aware sizing on startup",
               "ScrollPane with hidden scrollbars"]},
    # Slide 13 — Challenges
    {"type": "content", "title": "Challenges & Solutions",
     "items": ["##JavaFX MediaPlayer startup lag",
               "→ Replaced with mp3agic for metadata; no MediaPlayer on load",
               "##Scrollbar appearing in Library/Liked",
               "→ skinProperty listener hides bars after skin loads",
               "##Bundled JRE missing native DLLs",
               "→ Copy JavaFX bin/*.dll into runtime/bin/",
               "##Open With registration",
               "→ Inno Setup registry entries + FriendlyAppName key",
               "##Data not persisting across sessions",
               "→ Pipe-separated save file with artist & duration cached"]},
    # Slide 14 — Conclusion
    {"type": "content", "title": "Conclusion & Future Work",
     "items": ["##Achievements",
               "Full-featured music player with modern dark UI",
               "Bundled JRE installer — no Java required by user",
               "Persistent library with fast startup",
               "##Limitations",
               "Windows only (JavaFX Media needs GStreamer on Linux)",
               "No playlist creation yet",
               "##Future Work",
               "Equalizer & audio effects",
               "Online streaming integration",
               "Playlist & queue management",
               "Cross-platform Linux/macOS support"]},
    # Slide 15 — Thank You
    {"type": "thankyou"},
]

for i, data in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, BG_DARK)

    stype = data["type"]

    if stype == "title":
        # Big decorative circle top-right
        add_rect(slide, 9.5, -1.5, 5, 5, BG_MID)
        add_rect(slide, 10.5, -0.5, 3.5, 3.5, ACCENT)
        # Bottom strip
        add_rect(slide, 0, 6.8, 13.33, 0.7, ACCENT2)
        # Music note decoration
        add_text_box(slide, "♪", 10.2, 0.2, 2, 2, size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, "Auralin", 0.8, 1.5, 9, 1.5, size=72, bold=True, color=WHITE)
        add_text_box(slide, "Music Player", 0.8, 2.8, 9, 1.2, size=52, bold=True, color=ACCENT)
        divider_line(slide, 4.2)
        add_text_box(slide, "A Modern JavaFX Desktop Music Application", 0.8, 4.4, 10, 0.6, size=22, color=LIGHT_GRAY)
        add_text_box(slide, "Dip Karmokar  |  OOP Project  |  2026", 0.8, 5.2, 10, 0.5, size=18, color=LIGHT_GRAY, italic=True)

    elif stype == "thankyou":
        add_rect(slide, 0, 0, 13.33, 7.5, BG_MID)
        add_rect(slide, 0, 6.5, 13.33, 1.0, ACCENT2)
        add_rect(slide, 4.5, 1.0, 4.33, 4.5, ACCENT)
        add_text_box(slide, "♪", 5.5, 1.2, 2.5, 2.5, size=90, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, "Thank You!", 2.5, 3.8, 8.33, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, "Auralin Music Player", 3.5, 5.0, 6.33, 0.7, size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, "Dip Karmokar", 4.5, 5.7, 4.33, 0.5, size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    elif stype == "arch":
        add_rect(slide, 0, 0, 13.33, 1.0, BG_MID)
        add_text_box(slide, "System Architecture", 0.4, 0.1, 12, 0.8, size=32, bold=True, color=WHITE)
        divider_line(slide, 1.05)
        # MVC boxes
        boxes = [
            (0.4,  1.4, 3.5, 4.5, ACCENT2,  "MODEL",      ["Song.java", "title, artist", "filePath, liked", "plays, lastPlayed"]),
            (4.9,  1.4, 3.5, 4.5, ACCENT,   "VIEW",       ["HomeView", "LibraryView", "SongTable", "PlayerBar", "Sidebar"]),
            (9.4,  1.4, 3.5, 4.5, BG_MID,   "CONTROLLER", ["PlayerController", "ViewManager", "DatabaseService", "FileImportService"]),
        ]
        for bx, by, bw, bh, bc, bt, bi in boxes:
            add_rect(slide, bx, by, bw, bh, bc)
            add_text_box(slide, bt, bx+0.1, by+0.1, bw-0.2, 0.6, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            for j, item in enumerate(bi):
                add_text_box(slide, "• " + item, bx+0.2, by+0.8+j*0.7, bw-0.3, 0.6, size=15, color=WHITE)
        # Arrows
        add_text_box(slide, "⟷", 4.0, 3.3, 0.8, 0.5, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide, "⟷", 8.6, 3.3, 0.8, 0.5, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide, "Observer Pattern  |  MVC  |  Facade  |  Strategy", 1.5, 6.5, 10, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    else:
        # Standard content slide
        add_rect(slide, 0, 0, 13.33, 1.0, BG_MID)
        add_text_box(slide, data["title"], 0.4, 0.1, 12, 0.8, size=32, bold=True, color=WHITE)
        divider_line(slide, 1.05)
        # Decorative music note
        add_text_box(slide, "♪", 11.8, 0.05, 1.2, 0.9, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        # Slide number
        add_text_box(slide, str(i+1), 12.8, 6.9, 0.5, 0.4, size=12, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
        # Content
        items = data.get("items", [])
        mid = len(items) // 2
        add_bullets(slide, items[:mid], 0.4, 1.2, 6.0, 5.8)
        add_bullets(slide, items[mid:], 6.7, 1.2, 6.0, 5.8)
        # Bottom accent bar
        add_rect(slide, 0, 7.1, 13.33, 0.4, ACCENT2)

out = r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\AuralinPresentation.pptx"
prs.save(out)
print("Saved:", out)
