# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Minimal Premium Palette ───────────────────────────────────────────────
BLACK   = RGBColor(0x09, 0x09, 0x09)
DARK    = RGBColor(0x11, 0x11, 0x1A)
CARD    = RGBColor(0x1A, 0x1A, 0x28)
BORDER  = RGBColor(0x2A, 0x2A, 0x3A)
RED     = RGBColor(0xFA, 0x2D, 0x48)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x94, 0x94, 0xA8)
LGRAY   = RGBColor(0x3A, 0x3A, 0x50)

PICS = [
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (1).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (2).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (3).png",
    r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\Auralin picsss\ (4).png",
]

# ── Core helpers ──────────────────────────────────────────────────────────
def set_bg(slide, color=BLACK):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line_color=None, line_w=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def label(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, font="Calibri Light"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri" if bold else font
    return tb

def multiline(slide, lines, l, t, w, h, base=16):
    """lines: list of (text, size, bold, color, indent)"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (text, size, bold, color, indent) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = indent
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri" if bold else "Calibri Light"
    return tb

def pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return None

def no_advance(slide):
    """Manual click only — no auto advance."""
    old = slide._element.find(qn('p:transition'))
    if old is not None: slide._element.remove(old)
    tr = etree.SubElement(slide._element, qn('p:transition'))
    tr.set('spd', 'slow')
    etree.SubElement(tr, qn('p:fade'))

def anim_fade(slide, shape, delay=0, dur=500):
    """Fade-in on click."""
    try:
        timing = slide._element.find(qn('p:timing'))
        if timing is None:
            timing = etree.SubElement(slide._element, qn('p:timing'))
        tnLst = timing.find(qn('p:tnLst'))
        if tnLst is None:
            tnLst = etree.SubElement(timing, qn('p:tnLst'))
        par = tnLst.find(qn('p:par'))
        if par is None:
            par = etree.SubElement(tnLst, qn('p:par'))
            cTn = etree.SubElement(par, qn('p:cTn'))
            cTn.set('id','1'); cTn.set('dur','indefinite')
            cTn.set('restart','whenNotActive'); cTn.set('nodeType','tmRoot')
            cl = etree.SubElement(cTn, qn('p:childTnLst'))
            seq = etree.SubElement(cl, qn('p:seq'))
            seq.set('concurrent','1'); seq.set('nextAc','seek')
            cTn2 = etree.SubElement(seq, qn('p:cTn'))
            cTn2.set('id','2'); cTn2.set('dur','indefinite')
            cTn2.set('nodeType','mainSeq')
            etree.SubElement(cTn2, qn('p:childTnLst'))
            etree.SubElement(seq, qn('p:prevCondLst'))
            nc = etree.SubElement(seq, qn('p:nextCondLst'))
            cond = etree.SubElement(nc, qn('p:cond'))
            cond.set('evt','onNext'); cond.set('delay','0')
            tn = etree.SubElement(cond, qn('p:tn')); tn.set('val','2')
        seq = par.find(qn('p:cTn')).find(qn('p:childTnLst')).find(qn('p:seq'))
        ml = seq.find(qn('p:cTn')).find(qn('p:childTnLst'))
        idx = len(ml) + 3
        p2 = etree.SubElement(ml, qn('p:par'))
        c3 = etree.SubElement(p2, qn('p:cTn'))
        c3.set('id',str(idx)); c3.set('fill','hold')
        sc = etree.SubElement(c3, qn('p:stCondLst'))
        cd = etree.SubElement(sc, qn('p:cond')); cd.set('delay',str(delay))
        cl2 = etree.SubElement(c3, qn('p:childTnLst'))
        p3 = etree.SubElement(cl2, qn('p:par'))
        c4 = etree.SubElement(p3, qn('p:cTn'))
        c4.set('id',str(idx+1)); c4.set('fill','hold')
        sc2 = etree.SubElement(c4, qn('p:stCondLst'))
        cd2 = etree.SubElement(sc2, qn('p:cond')); cd2.set('delay','0')
        cl3 = etree.SubElement(c4, qn('p:childTnLst'))
        ae = etree.SubElement(cl3, qn('p:animEffect'))
        ae.set('transition','in'); ae.set('filter','fade')
        cb = etree.SubElement(ae, qn('p:cBhvr'))
        c5 = etree.SubElement(cb, qn('p:cTn'))
        c5.set('id',str(idx+2)); c5.set('dur',str(dur)); c5.set('fill','hold')
        te = etree.SubElement(cb, qn('p:tgtEl'))
        sp = etree.SubElement(te, qn('p:spTgt'))
        sp.set('spid',str(shape.shape_id))
    except: pass

# ── Reusable slide chrome ─────────────────────────────────────────────────
def slide_chrome(slide, accent=RED):
    set_bg(slide, BLACK)
    box(slide, 0, 0, 13.33, 7.5, DARK)
    # Left accent strip
    box(slide, 0, 0, 0.06, 7.5, accent)
    # Top bar
    box(slide, 0.06, 0, 13.27, 0.06, accent)
    # Bottom bar
    box(slide, 0, 7.44, 13.33, 0.06, accent)
    # Footer
    box(slide, 0.06, 7.1, 13.27, 0.34, CARD)
    label(slide, "Auralin Music Player  |  Dip Karmokar", 0.3, 7.14, 8, 0.26, size=11, color=GRAY)

def slide_header(slide, title, sub=None, accent=RED, page=None):
    t = label(slide, title, 0.3, 0.18, 11, 0.72, size=38, bold=True, color=WHITE)
    if sub:
        label(slide, sub, 0.3, 0.88, 11, 0.32, size=14, color=GRAY, italic=True)
    # Thin underline
    box(slide, 0.3, 1.22, 12.73, 0.04, accent)
    if page:
        label(slide, str(page), 12.6, 7.14, 0.6, 0.26, size=11, color=GRAY, align=PP_ALIGN.RIGHT)
    return t

# ═══════════════════════════════════════════════════════════════════════════
# S1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BLACK)
box(sl, 0, 0, 13.33, 7.5, DARK)
box(sl, 0, 0, 0.06, 7.5, RED)
box(sl, 0.06, 0, 13.27, 0.06, RED)
box(sl, 0, 7.44, 13.33, 0.06, RED)
# Big circle decoration right side
box(sl, 8.5, -0.5, 5.5, 5.5, CARD)
box(sl, 8.9, -0.1, 4.7, 4.7, BLACK)
box(sl, 9.3, 0.3, 3.9, 3.9, CARD)
box(sl, 9.7, 0.7, 3.1, 3.1, BLACK)
t_note = label(sl, "♪", 10.0, 0.8, 2.5, 2.5, size=96, bold=True, color=RED, align=PP_ALIGN.CENTER)
# Title block
box(sl, 0.3, 1.8, 0.08, 2.8, RED)
t1 = label(sl, "AURALIN", 0.55, 1.7, 9, 1.5, size=86, bold=True, color=WHITE)
t2 = label(sl, "MUSIC PLAYER", 0.55, 3.1, 9, 0.9, size=44, bold=True, color=RED)
box(sl, 0.55, 4.1, 6.5, 0.05, LGRAY)
t3 = label(sl, "A Modern JavaFX Desktop Music Application", 0.55, 4.25, 10, 0.5, size=19, color=GRAY, italic=True)
t4 = label(sl, "Dip Karmokar   |   OOP Project   |   2026", 0.55, 4.85, 10, 0.45, size=16, color=GRAY)
box(sl, 0, 6.6, 13.33, 0.84, CARD)
box(sl, 0, 6.6, 13.33, 0.05, LGRAY)
label(sl, "JavaFX 25   ·   mp3agic   ·   MVC Architecture   ·   Windows Installer   ·   Bundled JRE",
      0.3, 6.7, 12.73, 0.4, size=13, color=GRAY, align=PP_ALIGN.CENTER)
no_advance(sl)
for s,d in [(t1,0),(t2,150),(t3,300),(t4,450)]: anim_fade(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# S2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl, RED); slide_header(sl,"Agenda","What we'll cover today", RED, 2)
no_advance(sl)
items = [
    ("01  Introduction & Motivation", RED),
    ("02  System Architecture", PURPLE),
    ("03  Key Features", CYAN),
    ("04  Implementation & Patterns", RED),
    ("05  Data Persistence", PURPLE),
    ("06  UI/UX Design", CYAN),
    ("07  Challenges & Solutions", RED),
    ("08  Conclusion & Future Work", PURPLE),
]
for i,(text,col) in enumerate(items):
    row=i//2; ci=i%2
    fx=0.3+ci*6.5; fy=1.4+row*1.35
    r=box(sl,fx,fy,6.1,1.15,CARD)
    box(sl,fx,fy,0.06,1.15,col)
    t=label(sl,text,fx+0.2,fy+0.3,5.7,0.55,size=19,bold=True,color=WHITE)
    anim_fade(sl,r,i*100)

# ═══════════════════════════════════════════════════════════════════════════
# S3 — INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,RED); slide_header(sl,"Introduction","What is Auralin Music Player?",RED,3)
no_advance(sl)
c1=box(sl,0.3,1.35,6.1,5.5,CARD); c2=box(sl,6.8,1.35,6.2,5.5,CARD)
box(sl,0.3,1.35,6.1,0.05,RED); box(sl,6.8,1.35,6.2,0.05,PURPLE)
label(sl,"What is Auralin?",0.5,1.45,5.7,0.42,size=15,bold=True,color=RED)
label(sl,"Goals & Scope",7.0,1.45,5.8,0.42,size=15,bold=True,color=PURPLE)
left_lines = [
    ("Modern desktop music player",16,False,WHITE,0),
    ("Built with JavaFX 25 + Java",16,False,GRAY,0),
    ("Inspired by Spotify dark UI",16,False,GRAY,0),
    ("",12,False,GRAY,0),
    ("Motivation",15,True,RED,0),
    ("Existing players lack modern UI",15,False,GRAY,0),
    ("Deep dive into JavaFX MVC",15,False,GRAY,0),
    ("",12,False,GRAY,0),
    ("Tech Stack",15,True,RED,0),
    ("JavaFX 25, mp3agic, Inno Setup",15,False,GRAY,0),
    ("launch4j, jlink bundled JRE",15,False,GRAY,0),
]
right_lines = [
    ("Fast startup — no MediaPlayer on load",16,False,WHITE,0),
    ("Persistent library across sessions",16,False,GRAY,0),
    ("Open With — default MP3 player",16,False,GRAY,0),
    ("Custom frameless Win11 window",16,False,GRAY,0),
    ("",12,False,GRAY,0),
    ("Deliverables",15,True,PURPLE,0),
    ("Runnable .exe installer",15,False,GRAY,0),
    ("Bundled JRE — no Java needed",15,False,GRAY,0),
    ("Full MVC architecture",15,False,GRAY,0),
]
b1=multiline(sl,left_lines,0.5,1.95,5.7,4.8)
b2=multiline(sl,right_lines,7.0,1.95,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: anim_fade(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# S4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,PURPLE); slide_header(sl,"System Architecture","MVC + Observer + Facade + Strategy",PURPLE,4)
no_advance(sl)
arch=[
    (0.3,1.35,3.9,5.5,RED,"MODEL",["Song.java","title, artist","filePath","liked, plays","lastPlayed","getDurationStr()"]),
    (4.7,1.35,3.9,5.5,PURPLE,"VIEW",["HomeView","LibraryView","SearchView","LikedView","SongTable","PlayerBar","Sidebar","NowPlayingView"]),
    (9.1,1.35,3.9,5.5,CYAN,"CONTROLLER",["PlayerController","ViewManager","DatabaseService","FileImportService","PlayerBar.Listener","MusicCard.Listener"]),
]
for bx,by,bw,bh,bc,bt,bi in arch:
    r=box(sl,bx,by,bw,bh,CARD)
    box(sl,bx,by,bw,0.05,bc)
    label(sl,bt,bx+0.15,by+0.12,bw-0.3,0.48,size=17,bold=True,color=bc,align=PP_ALIGN.CENTER)
    box(sl,bx+0.15,by+0.65,bw-0.3,0.03,LGRAY)
    for j,item in enumerate(bi):
        label(sl,"  "+item,bx+0.15,by+0.75+j*0.72,bw-0.3,0.6,size=13,color=GRAY)
    anim_fade(sl,r,200)
label(sl,"<->",4.15,3.8,0.55,0.5,size=20,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
label(sl,"<->",8.6,3.8,0.55,0.5,size=20,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# S5 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,CYAN); slide_header(sl,"Key Features","Everything Auralin can do",CYAN,5)
no_advance(sl)
feats=[
    ("Library Management","Add files & folders, ID3 tags",RED),
    ("Real-time Search","Filter songs instantly",PURPLE),
    ("Liked Songs","Heart toggle + persistent",CYAN),
    ("Home View","Recent + Recommended grids",RED),
    ("Now Playing","Full-screen artwork overlay",PURPLE),
    ("Auto-Save","Library persists on restart",CYAN),
    ("Open With","Register as default MP3 player",RED),
    ("Custom Window","Frameless + Win11 controls",PURPLE),
]
pos=[(0.3,1.35),(3.55,1.35),(6.8,1.35),(10.05,1.35),
     (0.3,3.95),(3.55,3.95),(6.8,3.95),(10.05,3.95)]
for i,((fx,fy),(title,desc,col)) in enumerate(zip(pos,feats)):
    r=box(sl,fx,fy,3.0,2.3,CARD)
    box(sl,fx,fy,3.0,0.05,col)
    box(sl,fx,fy+2.25,3.0,0.05,col)
    label(sl,title,fx+0.15,fy+0.18,2.7,0.52,size=15,bold=True,color=col)
    label(sl,desc,fx+0.15,fy+0.82,2.7,0.65,size=13,color=GRAY)
    anim_fade(sl,r,i*80)

# ═══════════════════════════════════════════════════════════════════════════
# S6 — SCREENSHOTS 1 & 2
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,RED); slide_header(sl,"Home View & Library","Screenshots from Auralin Music Player",RED,6)
no_advance(sl)
box(sl,0.3,1.35,6.1,5.5,CARD); box(sl,6.8,1.35,6.2,5.5,CARD)
box(sl,0.3,1.35,6.1,0.05,RED); box(sl,6.8,1.35,6.2,0.05,PURPLE)
p1=pic(sl,PICS[0],0.4,1.45,5.9,4.1)
p2=pic(sl,PICS[1],6.9,1.45,5.9,4.1)
label(sl,"Home View — Recently Played & Recommendations",0.3,5.65,6.1,0.5,size=12,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
label(sl,"Library View — Song Table with Artwork",6.8,5.65,6.2,0.5,size=12,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
if p1: anim_fade(sl,p1,200)
if p2: anim_fade(sl,p2,400)

# ═══════════════════════════════════════════════════════════════════════════
# S7 — SCREENSHOTS 3 & 4
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,PURPLE); slide_header(sl,"Now Playing & Liked Songs","Full-screen overlay and liked songs view",PURPLE,7)
no_advance(sl)
box(sl,0.3,1.35,6.1,5.5,CARD); box(sl,6.8,1.35,6.2,5.5,CARD)
box(sl,0.3,1.35,6.1,0.05,CYAN); box(sl,6.8,1.35,6.2,0.05,RED)
p3=pic(sl,PICS[2],0.4,1.45,5.9,4.1)
p4=pic(sl,PICS[3],6.9,1.45,5.9,4.1)
label(sl,"Now Playing — Full-screen artwork overlay",0.3,5.65,6.1,0.5,size=12,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
label(sl,"Liked Songs — Filtered heart-toggled songs",6.8,5.65,6.2,0.5,size=12,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
if p3: anim_fade(sl,p3,200)
if p4: anim_fade(sl,p4,400)

# ═══════════════════════════════════════════════════════════════════════════
# S8 — PLAYBACK
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,RED); slide_header(sl,"Music Playback & Controls","PlayerController + PlayerBar",RED,8)
no_advance(sl)
c1=box(sl,0.3,1.35,6.1,5.5,CARD); c2=box(sl,6.8,1.35,6.2,5.5,CARD)
box(sl,0.3,1.35,6.1,0.05,RED); box(sl,6.8,1.35,6.2,0.05,PURPLE)
label(sl,"PlayerController",0.5,1.45,5.7,0.42,size=15,bold=True,color=RED)
label(sl,"PlayerBar Component",7.0,1.45,5.8,0.42,size=15,bold=True,color=PURPLE)
ll=[("JavaFX MediaPlayer for audio",15,False,WHITE,0),("Play / Pause / Next / Previous",15,False,GRAY,0),("Shuffle — random song selection",15,False,GRAY,0),("Repeat — loop current song",15,False,GRAY,0),("Seek bar — jump to position",15,False,GRAY,0),("Volume control slider",15,False,GRAY,0),("",12,False,GRAY,0),("Callbacks",14,True,RED,0),("onSongChange, onLikeChange",14,False,GRAY,0),("onPlayStateChange",14,False,GRAY,0)]
rl=[("Always visible at bottom",15,False,WHITE,0),("Shows artwork thumbnail",15,False,GRAY,0),("Song title + artist name",15,False,GRAY,0),("Progress bar with time",15,False,GRAY,0),("Heart button for like/unlike",15,False,GRAY,0),("Volume slider",15,False,GRAY,0),("",12,False,GRAY,0),("Design",14,True,PURPLE,0),("Listener interface pattern",14,False,GRAY,0),("Decoupled from controller",14,False,GRAY,0)]
b1=multiline(sl,ll,0.5,1.95,5.7,4.8); b2=multiline(sl,rl,7.0,1.95,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: anim_fade(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# S9 — DESIGN PATTERNS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,PURPLE); slide_header(sl,"Design Patterns Used","Clean architecture through proven patterns",PURPLE,9)
no_advance(sl)
patterns=[
    ("MVC","Model-View-Controller separates data, UI, and logic cleanly",RED),
    ("Observer","Callbacks: onSongChange, onLikeChange, onPlayStateChange",PURPLE),
    ("Strategy","Shuffle vs Sequential playback — swappable at runtime",CYAN),
    ("Facade","DatabaseService hides mp3agic + MediaPlayer complexity",RED),
    ("Factory","SongTableListener / MusicCardListener interface factories",PURPLE),
    ("Singleton","DatabaseService — single source of truth for all song data",CYAN),
]
for i,(name,desc,col) in enumerate(patterns):
    row=i//2; ci=i%2
    fx=0.3+ci*6.5; fy=1.35+row*1.85
    r=box(sl,fx,fy,6.2,1.65,CARD)
    box(sl,fx,fy,0.06,1.65,col)
    box(sl,fx,fy,6.2,0.05,col)
    label(sl,name,fx+0.2,fy+0.15,2.5,0.52,size=19,bold=True,color=col)
    label(sl,desc,fx+0.2,fy+0.82,5.8,0.65,size=13,color=GRAY)
    anim_fade(sl,r,i*120)

# ═══════════════════════════════════════════════════════════════════════════
# S10 — DATA PERSISTENCE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,CYAN); slide_header(sl,"Data Persistence & Storage","Library survives app restarts",CYAN,10)
no_advance(sl)
r0=box(sl,0.3,1.35,12.73,1.0,CARD)
box(sl,0.3,1.35,12.73,0.05,CYAN)
label(sl,"Save File:  %APPDATA%\\AuralinPlayer\\library.dat",0.5,1.42,12.3,0.38,size=14,bold=True,color=WHITE)
label(sl,"filePath  |  liked  |  plays  |  lastPlayed  |  artist  |  durationMs",0.5,1.8,12.3,0.38,size=13,color=CYAN,font="Courier New")
c1=box(sl,0.3,2.5,6.1,4.35,CARD); c2=box(sl,6.8,2.5,6.2,4.35,CARD)
box(sl,0.3,2.5,6.1,0.05,RED); box(sl,6.8,2.5,6.2,0.05,PURPLE)
label(sl,"Save Triggers",0.5,2.58,5.7,0.42,size=15,bold=True,color=RED)
label(sl,"Load on Startup",7.0,2.58,5.8,0.42,size=15,bold=True,color=PURPLE)
ll=[("On song play — recordPlay()",15,False,WHITE,0),("On like toggle — saveToDisk()",15,False,GRAY,0),("On app close — stop()",15,False,GRAY,0),("Background thread — non-blocking",15,False,GRAY,0),("",12,False,GRAY,0),("Why Fast?",14,True,RED,0),("mp3agic reads tags directly",14,False,GRAY,0),("No MediaPlayer on startup",14,False,GRAY,0),("Thread pool (3 threads)",14,False,GRAY,0)]
rl=[("Read pipe-separated lines",15,False,WHITE,0),("Skip missing files automatically",15,False,GRAY,0),("Restore: liked, plays, lastPlayed",15,False,GRAY,0),("Restore: artist, durationMs",15,False,GRAY,0),("",12,False,GRAY,0),("Artwork Loading",14,True,PURPLE,0),("Background thread per song",14,False,GRAY,0),("Platform.runLater for UI update",14,False,GRAY,0),("Table refreshes incrementally",14,False,GRAY,0)]
b1=multiline(sl,ll,0.5,3.05,5.8,3.6); b2=multiline(sl,rl,7.0,3.05,5.9,3.6)
for s,d in [(r0,0),(c1,100),(c2,200),(b1,300),(b2,400)]: anim_fade(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# S11 — UI/UX
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,RED); slide_header(sl,"UI/UX Design Highlights","Premium dark aesthetic inspired by Spotify",RED,11)
no_advance(sl)
palette=[("#090909","Background"),("#1A1A28","Card"),("#FA2D48","Auralin Red"),("#8B5CF6","Purple"),("#06B6D4","Cyan"),("#94A3B8","Gray")]
for i,(hx,nm) in enumerate(palette):
    fx=0.3+i*2.15
    r=box(sl,fx,1.35,1.9,1.1,RGBColor(int(hx[1:3],16),int(hx[3:5],16),int(hx[5:7],16)))
    box(sl,fx,1.35,1.9,0.05,WHITE)
    label(sl,nm,fx,2.52,1.9,0.38,size=12,color=GRAY,align=PP_ALIGN.CENTER)
    label(sl,hx,fx,2.88,1.9,0.3,size=10,color=LGRAY,align=PP_ALIGN.CENTER,italic=True)
c1=box(sl,0.3,3.25,6.1,3.6,CARD); c2=box(sl,6.8,3.25,6.2,3.6,CARD)
box(sl,0.3,3.25,6.1,0.05,RED); box(sl,6.8,3.25,6.2,0.05,PURPLE)
label(sl,"Custom Window",0.5,3.33,5.7,0.42,size=15,bold=True,color=RED)
label(sl,"Components",7.0,3.33,5.8,0.42,size=15,bold=True,color=PURPLE)
multiline(sl,[("StageStyle.TRANSPARENT",14,False,WHITE,0),("Custom title bar (32px)",14,False,GRAY,0),("Win11-style min/max/close",14,False,GRAY,0),("8-edge resize handles",14,False,GRAY,0),("Rounded corners (radius 10)",14,False,GRAY,0),("Screen-aware sizing on start",14,False,GRAY,0)],0.5,3.78,5.7,2.9)
multiline(sl,[("MusicCard — 160px hover card",14,False,WHITE,0),("SongTable — custom TableView",14,False,GRAY,0),("PlayerBar — bottom fixed bar",14,False,GRAY,0),("Sidebar — navigation panel",14,False,GRAY,0),("NowPlayingView — overlay",14,False,GRAY,0),("ScrollPane hidden scrollbars",14,False,GRAY,0)],7.0,3.78,5.9,2.9)

# ═══════════════════════════════════════════════════════════════════════════
# S12 — CHALLENGES
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,PURPLE); slide_header(sl,"Challenges & Solutions","Real problems solved during development",PURPLE,12)
no_advance(sl)
chs=[("Slow startup","MediaPlayer for each song on load","mp3agic reads tags; no MediaPlayer needed",RED),("Scrollbar appearing","TableView internal ScrollPane","skinProperty listener hides bars post-skin",PURPLE),("Missing DLLs in EXE","jlink skips JavaFX native DLLs","Copy javafx/bin/*.dll into runtime/bin/",CYAN),("Data not persisting","No save mechanism existed","Pipe-separated file with all metadata cached",RED),("Open With missing","App not registered in Windows","Inno Setup registry + FriendlyAppName key",PURPLE)]
for i,(prob,cause,sol,col) in enumerate(chs):
    fy=1.35+i*1.1
    r=box(sl,0.3,fy,12.73,0.98,CARD)
    box(sl,0.3,fy,0.06,0.98,col)
    box(sl,0.3,fy,12.73,0.04,col)
    label(sl,prob,0.5,fy+0.1,2.8,0.4,size=14,bold=True,color=col)
    label(sl,"Cause: "+cause,3.5,fy+0.1,4.5,0.36,size=12,color=GRAY,italic=True)
    label(sl,"Fix: "+sol,8.2,fy+0.1,4.7,0.36,size=12,color=WHITE)
    anim_fade(sl,r,i*150)

# ═══════════════════════════════════════════════════════════════════════════
# S13 — INSTALLER
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,CYAN); slide_header(sl,"Installer & Deployment","Zero-dependency Windows installer",CYAN,13)
no_advance(sl)
steps=[("1","Compile Java","javac -> .class",RED),("2","Fat JAR","jar + mp3agic",PURPLE),("3","jlink JRE","80MB bundled JRE",CYAN),("4","launch4j","JAR -> .exe",RED),("5","Inno Setup","Full installer",PURPLE)]
for i,(num,title,desc,col) in enumerate(steps):
    fx=0.3+i*2.6; fy=1.4
    r=box(sl,fx,fy,2.4,2.2,CARD)
    box(sl,fx,fy,2.4,0.05,col)
    box(sl,fx,fy+2.15,2.4,0.05,col)
    label(sl,num,fx+0.1,fy+0.12,0.5,0.5,size=22,bold=True,color=col)
    label(sl,title,fx+0.1,fy+0.65,2.2,0.48,size=15,bold=True,color=WHITE)
    label(sl,desc,fx+0.1,fy+1.18,2.2,0.75,size=12,color=GRAY)
    if i<4: label(sl,"->",fx+2.4,fy+0.85,0.2,0.4,size=18,bold=True,color=LGRAY,align=PP_ALIGN.CENTER)
    anim_fade(sl,r,i*180)
r2=box(sl,0.3,3.8,12.73,3.0,CARD)
box(sl,0.3,3.8,12.73,0.05,CYAN)
label(sl,"What the installer does",0.5,3.88,12.3,0.42,size=15,bold=True,color=CYAN)
multiline(sl,[("Installs AuralinPlayer.exe + bundled runtime to Program Files",14,False,WHITE,0),("Registers .mp3 Open With association in Windows registry",14,False,GRAY,0),("Creates Desktop + Start Menu shortcuts with custom icon",14,False,GRAY,0),("Includes EULA acceptance screen  |  Full uninstall support",14,False,GRAY,0)],0.5,4.38,12.3,2.2)
anim_fade(sl,r2,900)

# ═══════════════════════════════════════════════════════════════════════════
# S14 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_chrome(sl,RED); slide_header(sl,"Conclusion & Future Work","What was achieved and what comes next",RED,14)
no_advance(sl)
c1=box(sl,0.3,1.35,6.1,5.5,CARD); c2=box(sl,6.8,1.35,6.2,5.5,CARD)
box(sl,0.3,1.35,6.1,0.05,RED); box(sl,6.8,1.35,6.2,0.05,PURPLE)
label(sl,"Achievements",0.5,1.45,5.7,0.42,size=15,bold=True,color=RED)
label(sl,"Future Work",7.0,1.45,5.8,0.42,size=15,bold=True,color=PURPLE)
ll=[("Full-featured music player",15,False,WHITE,0),("Modern Spotify-inspired dark UI",15,False,GRAY,0),("MVC architecture with patterns",15,False,GRAY,0),("Persistent library — fast startup",15,False,GRAY,0),("Bundled JRE installer",15,False,GRAY,0),("Open With MP3 registration",15,False,GRAY,0),("Custom frameless Win11 window",15,False,GRAY,0),("Background metadata loading",15,False,GRAY,0)]
rl=[("Equalizer & audio effects",15,False,WHITE,0),("Online streaming integration",15,False,GRAY,0),("Playlist & queue management",15,False,GRAY,0),("Cross-platform Linux/macOS",15,False,GRAY,0),("Album & artist grouping",15,False,GRAY,0),("Lyrics display integration",15,False,GRAY,0),("Mini player mode",15,False,GRAY,0),("Cloud sync for library",15,False,GRAY,0)]
b1=multiline(sl,ll,0.5,1.95,5.7,4.8); b2=multiline(sl,rl,7.0,1.95,5.9,4.8)
for s,d in [(c1,0),(c2,100),(b1,200),(b2,300)]: anim_fade(sl,s,d)

# ═══════════════════════════════════════════════════════════════════════════
# S15 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl,BLACK)
box(sl,0,0,13.33,7.5,DARK)
box(sl,0,0,0.06,7.5,RED)
box(sl,0.06,0,13.27,0.06,RED)
box(sl,0,7.44,13.33,0.06,RED)
# Center card
box(sl,3.5,0.9,6.33,5.7,CARD)
box(sl,3.5,0.9,6.33,0.06,RED)
box(sl,3.5,6.54,6.33,0.06,PURPLE)
box(sl,3.5,0.9,0.06,5.7,RED)
box(sl,9.77,0.9,0.06,5.7,PURPLE)
t1=label(sl,"♪",5.5,1.1,2.5,2.5,size=96,bold=True,color=RED,align=PP_ALIGN.CENTER)
t2=label(sl,"Thank You!",3.7,3.7,5.93,1.1,size=52,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
box(sl,4.5,4.85,4.33,0.04,LGRAY)
t3=label(sl,"Auralin Music Player",4.0,4.98,5.33,0.52,size=18,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
t4=label(sl,"Dip Karmokar   |   OOP Project   |   2026",3.7,5.6,5.93,0.42,size=14,color=GRAY,align=PP_ALIGN.CENTER)
box(sl,0,6.6,13.33,0.84,CARD)
box(sl,0,6.6,13.33,0.05,LGRAY)
label(sl,"Questions? Let's discuss!",0.3,6.7,12.73,0.4,size=14,color=GRAY,align=PP_ALIGN.CENTER,italic=True)
no_advance(sl)
for s,d in [(t1,0),(t2,200),(t3,400),(t4,600)]: anim_fade(sl,s,d)

# ── Save ──────────────────────────────────────────────────────────────────
out = r"C:\Users\Dip karmokar\IdeaProjects\Projectforoop\AuralinPresentation.pptx"
prs.save(out)
print("Done:", out)
